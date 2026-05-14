#!/usr/bin/env python3
"""
Antuario Brandbook — Pipeline de render.

1. Lee brandbook.html.
2. Extrae <head> + SVG defs y cada <section class="slide ..."> a archivos individuales.
3. Renderiza cada slide con Chrome headless a PNG @3x.
4. Arma PPTX 16:9 con una imagen por slide.
5. Ensambla PDF directo desde los PNGs (con PIL, calidad lossless).
"""

import os
import re
import shutil
import subprocess
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent
BB_DIR = ROOT.parent
HTML_PATH = BB_DIR / "brandbook.html"
SLIDES_DIR = ROOT / "slides"
PNG_DIR = ROOT / "png"
OUT_PPTX = BB_DIR / "Antuario-Brandbook.pptx"
OUT_PDF = BB_DIR / "Antuario-Brandbook.pdf"

CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"

# Slide dimensions in mm (from @page in CSS): 338.67 x 190.5 (16:9 landscape).
# 1mm = 3.7795275591 px @96dpi → 1280 x 720 px @1x.
# We render @3x → 3840 x 2160 (4K). Crisp for both PPTX & PDF.
SCALE = 3
WIDTH_PX_1X = 1280
HEIGHT_PX_1X = 720
WIDTH_PX = WIDTH_PX_1X * SCALE
HEIGHT_PX = HEIGHT_PX_1X * SCALE


def read_html() -> str:
    return HTML_PATH.read_text(encoding="utf-8")


def split_slides(html: str):
    """Return (head_html, [section_html, ...])."""
    head_match = re.search(r"(<!DOCTYPE.*?</head>)", html, flags=re.DOTALL | re.IGNORECASE)
    if not head_match:
        raise SystemExit("Could not find <head>")
    head = head_match.group(1)

    # Extract the SVG <defs> block that lives at top of <body>.
    svg_defs_match = re.search(
        r"(<svg width=\"0\" height=\"0\".*?</svg>)",
        html,
        flags=re.DOTALL,
    )
    svg_defs = svg_defs_match.group(1) if svg_defs_match else ""

    # All <section class="slide ...">...</section>
    sections = re.findall(
        r"(<section class=\"slide[^\"]*\".*?</section>)",
        html,
        flags=re.DOTALL,
    )
    return head, svg_defs, sections


def build_slide_doc(head: str, svg_defs: str, section: str) -> str:
    """Wrap a single section into a standalone HTML doc sized to the slide.

    Critical: when rendered to a viewport of WIDTH_PX_1X x HEIGHT_PX_1X (with
    device-scale-factor SCALE), we want the .slide to fill the viewport.

    The CSS already sizes .slide as 338.67mm x 190.5mm and the @page is also
    that size. Body has 0 margin. So a viewport of the exact px size at 96dpi
    will fit perfectly.
    """
    # Inject an override so the slide always renders flat at the viewport size
    # without page-break behavior or scrolling.
    override = """
<style>
  html, body { width: 338.67mm; height: 190.5mm; overflow: hidden; }
  body { margin: 0; padding: 0; }
  .slide { page-break-after: auto !important; }
</style>
"""
    return f"{head}\n{override}\n<body>\n{svg_defs}\n{section}\n</body>\n</html>"


def render_slide(slide_html_path: Path, png_path: Path):
    """Use headless Chrome to capture a PNG."""
    cmd = [
        CHROME,
        "--headless=new",
        "--disable-gpu",
        "--hide-scrollbars",
        "--no-sandbox",
        "--default-background-color=00000000",
        f"--force-device-scale-factor={SCALE}",
        f"--window-size={WIDTH_PX_1X},{HEIGHT_PX_1X}",
        f"--screenshot={png_path}",
        f"--virtual-time-budget=2000",
        slide_html_path.as_uri(),
    ]
    res = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
    if not png_path.exists():
        print("STDOUT:", res.stdout)
        print("STDERR:", res.stderr)
        raise SystemExit(f"Chrome didn't produce {png_path}")


def render_all():
    html = read_html()
    head, svg_defs, sections = split_slides(html)
    print(f"Found {len(sections)} slides.")

    if SLIDES_DIR.exists():
        shutil.rmtree(SLIDES_DIR)
    SLIDES_DIR.mkdir(parents=True)
    if PNG_DIR.exists():
        shutil.rmtree(PNG_DIR)
    PNG_DIR.mkdir(parents=True)

    png_paths = []
    for i, section in enumerate(sections, start=1):
        slide_html = build_slide_doc(head, svg_defs, section)
        slide_path = SLIDES_DIR / f"slide_{i:02d}.html"
        slide_path.write_text(slide_html, encoding="utf-8")

        png_path = PNG_DIR / f"slide_{i:02d}.png"
        print(f"  → Rendering slide {i:02d} ...", flush=True)
        render_slide(slide_path, png_path)
        png_paths.append(png_path)

    return png_paths


def build_pptx(png_paths):
    """Build a PPTX where each slide is exactly one full-bleed image."""
    from pptx import Presentation
    from pptx.util import Mm

    prs = Presentation()
    # Set slide size to match HTML: 338.67mm x 190.5mm (custom 16:9).
    prs.slide_width = Mm(338.67)
    prs.slide_height = Mm(190.5)

    blank_layout = prs.slide_layouts[6]  # blank

    for png in png_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(png),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(OUT_PPTX)
    print(f"PPTX written: {OUT_PPTX}")


def build_pdf(png_paths):
    """Build the PDF directly from PNGs. Each PNG = one PDF page."""
    from PIL import Image

    images = []
    for p in png_paths:
        im = Image.open(p)
        # PDF doesn't support alpha — flatten on white to be safe.
        if im.mode != "RGB":
            bg = Image.new("RGB", im.size, (255, 255, 255))
            if im.mode in ("RGBA", "LA"):
                bg.paste(im, mask=im.split()[-1])
            else:
                bg.paste(im)
            im = bg
        images.append(im)

    first, rest = images[0], images[1:]
    # 200 dpi makes the page size in PDF reasonable; image embeds losslessly.
    first.save(
        OUT_PDF,
        save_all=True,
        append_images=rest,
        format="PDF",
        resolution=288.0,  # matches our 3x render
    )
    print(f"PDF written: {OUT_PDF}")


def main():
    if not Path(CHROME).exists():
        raise SystemExit(f"Chrome not found at {CHROME}")
    pngs = render_all()
    build_pptx(pngs)
    build_pdf(pngs)
    print("\nDone.")
    print(f"  PPTX: {OUT_PPTX}")
    print(f"  PDF:  {OUT_PDF}")


if __name__ == "__main__":
    main()
